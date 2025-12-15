# -*- coding: utf-8 -*-
"""
속리산 축제 PPT 제안서 - 웹페이지 스크린샷 기반 고품질 PPT 생성
"""

import asyncio
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches
import os
import shutil

BASE_URL = "https://songnisan-festival-2026.vercel.app"
SCREENSHOT_DIR = r"C:\Users\user\songnisan-festival-2026\screenshots"
OUTPUT_PPT = r"C:\Users\user\songnisan-festival-2026\속리산축제_제안서.pptx"

# 슬라이드 구성: (URL 경로, 스크롤 Y 위치, 대기 시간)
SLIDES_CONFIG = [
    # 슬라이드 1: 홈 히어로 (표지)
    ("/", 0, 2),
    # 슬라이드 2: Executive Summary
    ("/", 750, 1),
    # 슬라이드 3: 배경 - 문제점/해결책
    ("/background", 0, 2),
    # 슬라이드 4: 배경 - 경쟁축제 벤치마킹
    ("/background", 900, 1),
    # 슬라이드 5: 여름축제 헤더
    ("/summer", 0, 2),
    # 슬라이드 6: 여름축제 프로그램
    ("/summer", 950, 1),
    # 슬라이드 7: 겨울축제 헤더
    ("/winter", 0, 2),
    # 슬라이드 8: 겨울축제 프로그램
    ("/winter", 950, 1),
    # 슬라이드 9: 실현가능성 - 예산
    ("/feasibility", 0, 2),
    # 슬라이드 10: 실현가능성 - 조직도/타임라인
    ("/feasibility", 1400, 1),
    # 슬라이드 11: 마케팅 전략
    ("/marketing", 0, 2),
    # 슬라이드 12: 기대효과
    ("/impact", 0, 2),
]


async def take_screenshots():
    """Playwright로 고품질 웹페이지 스크린샷 촬영"""
    # 기존 스크린샷 폴더 초기화
    if os.path.exists(SCREENSHOT_DIR):
        shutil.rmtree(SCREENSHOT_DIR)
    os.makedirs(SCREENSHOT_DIR)

    screenshot_paths = []

    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        context = await browser.new_context(
            viewport={"width": 1920, "height": 1080},
            device_scale_factor=2  # 4K 품질 (3840x2160)
        )
        page = await context.new_page()

        for i, (path, scroll_y, wait_time) in enumerate(SLIDES_CONFIG, 1):
            url = f"{BASE_URL}{path}"
            print(f"[{i}/12] Capturing: {path} (scroll: {scroll_y}px)")

            # 페이지 로드
            await page.goto(url, wait_until="networkidle")
            await asyncio.sleep(wait_time)  # 애니메이션 완료 대기

            # 스크롤
            if scroll_y > 0:
                await page.evaluate(f"window.scrollTo(0, {scroll_y})")
                await asyncio.sleep(0.8)

            # 스크린샷 촬영
            screenshot_path = os.path.join(SCREENSHOT_DIR, f"slide_{i:02d}.png")
            await page.screenshot(path=screenshot_path, full_page=False)
            screenshot_paths.append(screenshot_path)
            print(f"       Saved: slide_{i:02d}.png")

        await browser.close()

    return screenshot_paths


def create_ppt(screenshot_paths):
    """스크린샷으로 PPT 생성"""
    prs = Presentation()
    # 16:9 비율 설정
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # 빈 슬라이드 레이아웃
    blank_layout = prs.slide_layouts[6]

    for i, screenshot_path in enumerate(screenshot_paths, 1):
        print(f"[{i}/12] Creating slide {i}")
        slide = prs.slides.add_slide(blank_layout)

        # 스크린샷을 슬라이드 전체에 배치 (여백 없음)
        slide.shapes.add_picture(
            screenshot_path,
            Inches(0), Inches(0),
            width=Inches(16), height=Inches(9)
        )

    # PPT 저장
    prs.save(OUTPUT_PPT)
    print(f"\nPPT saved: {OUTPUT_PPT}")


async def main():
    print("=" * 50)
    print("Songnisan Festival PPT Generator")
    print("=" * 50)

    print("\n[Step 1] Taking high-quality screenshots...")
    screenshot_paths = await take_screenshots()

    print("\n[Step 2] Creating PPT...")
    create_ppt(screenshot_paths)

    print("\nDone!")


if __name__ == "__main__":
    asyncio.run(main())
