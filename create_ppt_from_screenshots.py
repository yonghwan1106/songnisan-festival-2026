"""
속리산 축제 인터랙티브 웹페이지 스크린샷 촬영 및 PPT 생성 스크립트
"""

import asyncio
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BASE_URL = "https://songnisan-festival-2026.vercel.app"
SCREENSHOT_DIR = r"C:\Users\user\songnisan-festival-2026\screenshots"
OUTPUT_PPT = r"C:\Users\user\songnisan-festival-2026\속리산축제_제안서_웹스크린샷.pptx"

# 슬라이드 구성: (슬라이드 제목, URL 경로, 스크롤 위치 또는 None)
SLIDES_CONFIG = [
    ("표지 - 속리산 사계 페스티벌", "/", 0),
    ("Executive Summary", "/", 800),
    ("기획 배경 - 계절성 문제", "/background", 0),
    ("경쟁 축제 벤치마킹", "/background", 1200),
    ("여름축제 - 호러 트레킹", "/summer", 0),
    ("여름축제 - 워터밤 & 피크닉", "/summer", 1200),
    ("겨울축제 - 미디어아트", "/winter", 0),
    ("겨울축제 - 소원등 & 푸드존", "/winter", 1200),
    ("예산 계획", "/feasibility", 0),
    ("조직도 & 추진일정", "/feasibility", 1500),
    ("마케팅 전략", "/marketing", 0),
    ("기대효과 & 심사기준 대응", "/impact", 0),
]


async def take_screenshots():
    """Playwright로 웹페이지 스크린샷 촬영"""
    os.makedirs(SCREENSHOT_DIR, exist_ok=True)

    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        context = await browser.new_context(
            viewport={"width": 1920, "height": 1080},
            device_scale_factor=1.5  # 고해상도
        )
        page = await context.new_page()

        screenshot_paths = []

        for i, (title, path, scroll_y) in enumerate(SLIDES_CONFIG, 1):
            url = f"{BASE_URL}{path}"
            print(f"[{i}/{len(SLIDES_CONFIG)}] 촬영 중: {title}")

            await page.goto(url, wait_until="networkidle")
            await asyncio.sleep(1)  # 애니메이션 완료 대기

            if scroll_y > 0:
                await page.evaluate(f"window.scrollTo(0, {scroll_y})")
                await asyncio.sleep(0.5)

            screenshot_path = os.path.join(SCREENSHOT_DIR, f"slide_{i:02d}.png")
            await page.screenshot(path=screenshot_path, full_page=False)
            screenshot_paths.append(screenshot_path)
            print(f"   저장: {screenshot_path}")

        await browser.close()

    return screenshot_paths


def create_ppt(screenshot_paths):
    """python-pptx로 PPT 생성"""
    prs = Presentation()
    prs.slide_width = Inches(16)  # 16:9 비율
    prs.slide_height = Inches(9)

    # 빈 슬라이드 레이아웃
    blank_layout = prs.slide_layouts[6]

    for i, (screenshot_path, (title, _, _)) in enumerate(zip(screenshot_paths, SLIDES_CONFIG), 1):
        print(f"[{i}/{len(screenshot_paths)}] 슬라이드 생성: {title}")

        slide = prs.slides.add_slide(blank_layout)

        # 스크린샷을 슬라이드 전체에 배치
        slide.shapes.add_picture(
            screenshot_path,
            Inches(0), Inches(0),
            width=Inches(16), height=Inches(9)
        )

        # 슬라이드 제목 오버레이 (반투명 배경)
        # 상단 좌측에 제목 표시
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.2),
            Inches(6), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i}. {title}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

    # PPT 저장
    prs.save(OUTPUT_PPT)
    print(f"\nPPT 저장 완료: {OUTPUT_PPT}")


async def main():
    print("=" * 60)
    print("속리산 축제 제안서 PPT 생성 스크립트")
    print("=" * 60)

    print("\n[1단계] 웹페이지 스크린샷 촬영...")
    screenshot_paths = await take_screenshots()

    print("\n[2단계] PPT 생성...")
    create_ppt(screenshot_paths)

    print("\n완료!")
    print(f"- 스크린샷: {SCREENSHOT_DIR}")
    print(f"- PPT 파일: {OUTPUT_PPT}")


if __name__ == "__main__":
    asyncio.run(main())
