# -*- coding: utf-8 -*-
"""
2026 속리산축제 콘텐츠 공모전 제안서 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 색상 정의
GREEN = RGBColor(0x2E, 0x7D, 0x32)  # 속리산 녹색
BLUE = RGBColor(0x15, 0x65, 0xC0)   # 겨울 파랑
ORANGE = RGBColor(0xE6, 0x51, 0x00) # 여름 주황
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)

def add_title_slide(prs, title, subtitle=""):
    """표지 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
    slide = prs.slides.add_slide(slide_layout)

    # 배경색 (녹색 그라데이션 효과는 단색으로 대체)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = GREEN
    background.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 부제목
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # 하단 정보
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026 속리산축제 콘텐츠 공모전"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_items, accent_color=GREEN):
    """일반 콘텐츠 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 상단 색상 바
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent_color
    top_bar.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK

    # 본문
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.space_after = Pt(12)

    return slide

def add_table_slide(prs, title, headers, rows, accent_color=GREEN):
    """표가 포함된 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 상단 색상 바
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent_color
    top_bar.line.fill.background()

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK

    # 표 생성
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table_width = Inches(9)
    table_height = Inches(0.4) * num_rows
    left = Inches(0.5)
    top = Inches(1.3)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table

    # 헤더 스타일
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent_color
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    # 데이터 행
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = DARK
            para.alignment = PP_ALIGN.CENTER

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, left_color=ORANGE, right_color=BLUE):
    """2단 레이아웃 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 상단 색상 바
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GREEN
    top_bar.line.fill.background()

    # 메인 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK

    # 왼쪽 컬럼 박스
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.3), Inches(4.5), Inches(5.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = left_color
    left_box.line.fill.background()

    # 왼쪽 제목
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.1), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 왼쪽 내용
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(4.1), Inches(4.5))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)

    # 오른쪽 컬럼 박스
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.3), Inches(4.5), Inches(5.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = right_color
    right_box.line.fill.background()

    # 오른쪽 제목
    right_title_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.5), Inches(4.1), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 오른쪽 내용
    right_content = slide.shapes.add_textbox(Inches(5.4), Inches(2.1), Inches(4.1), Inches(4.5))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_after = Pt(8)

    return slide

def create_presentation():
    """메인 함수: 12페이지 PPT 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== 슬라이드 1: 표지 ==========
    add_title_slide(
        prs,
        "2026 속리산축제\n콘텐츠 공모전 제안서",
        "기후 위기 시대, 신화와 자연의 재해석"
    )

    # ========== 슬라이드 2: Executive Summary ==========
    add_table_slide(
        prs,
        "Executive Summary",
        ["구분", "여름 축제", "겨울 축제"],
        [
            ["축제명", "속리산 야행: 도깨비와 걷는 숲", "속리산 겨울연화: 빛으로 피어나는 산사"],
            ["시기", "2026년 6월 27일 ~ 7월 12일", "2026년 12월 19일 ~ 2027년 1월 3일"],
            ["핵심 콘텐츠", "호러 트레킹, 워터밤, 에코 피크닉", "법주사 미디어아트, 소원등 트레킹"],
            ["타겟층", "MZ세대, 가족 단위", "가족, 연인, 중장년 힐링족"],
            ["예상 방문객", "8만명 (일 5,000명)", "6.4만명 (일 4,000명)"],
            ["예산", "8억 원", "10억 원"],
        ]
    )

    # ========== 슬라이드 3: 배경 및 목적 ==========
    add_content_slide(
        prs,
        "왜 여름·겨울 축제가 필요한가?",
        [
            "현재 보은군 관광의 문제점",
            "",
            "  - 봄(신화여행축제), 가을(대추축제)에 편중",
            "  - 여름·겨울 비수기 → 지역 상권 매출 절벽",
            "  - 기후 변화로 전통 겨울 축제(얼음낚시) 위협",
            "",
            "해결책: 계절성(Seasonality) 극복",
            "",
            "  - 여름: 자연의 냉기 + 심리적 냉기(공포) = Cool & Horror",
            "  - 겨울: 빛과 문화유산 = 기후 독립형 콘텐츠",
            "",
            "목표: 사계절 365일 체류형 관광 허브"
        ]
    )

    # ========== 슬라이드 4: 여름 축제 개요 ==========
    add_content_slide(
        prs,
        "여름 축제: 속리산 야행(夜行)",
        [
            "\"천년의 숲, 깨어나지 말아야 할 것이 깨어났다\"",
            "",
            "축제 개요",
            "  - 명칭: 2026 속리산 쿨(Cool) & 칠(Chill) 페스티벌",
            "  - 기간: 2026년 6월 27일(토) ~ 7월 12일(일), 16일간",
            "  - 장소: 오리숲, 세조길, 속리산 잔디광장",
            "",
            "핵심 가치",
            "  - 야간 관광 활성화",
            "  - MZ세대 청년 유입",
            "  - 친환경 축제 (Zero Waste)"
        ],
        accent_color=ORANGE
    )

    # ========== 슬라이드 5: 여름 축제 콘텐츠 ==========
    add_content_slide(
        prs,
        "여름 축제 킬러 콘텐츠",
        [
            "1. 오리숲 호러 트레킹 '도깨비의 금기'",
            "   - 국내 유일 국립공원 야간 호러 체험",
            "   - 도깨비 설화 기반 인터랙티브 미션",
            "   - 5개의 봉인 부적을 찾아 도깨비를 봉인하라!",
            "",
            "2. 속리산 워터밤 '대추나무 물총 전쟁'",
            "   - 대추 수호대 vs 도깨비 군단 물총 대결",
            "   - 쿨링 포그 & 인공 강우 존 운영",
            "",
            "3. 계곡 에코 피크닉 '신선놀음'",
            "   - 프라이빗 피크닉 존 예약제",
            "   - 보은 대추 도시락 + 어쿠스틱 버스킹"
        ],
        accent_color=ORANGE
    )

    # ========== 슬라이드 6: 겨울 축제 개요 ==========
    add_content_slide(
        prs,
        "겨울 축제: 속리산 겨울연화(蓮花)",
        [
            "\"천년의 빛, 당신의 소원을 비추다\"",
            "",
            "축제 개요",
            "  - 명칭: 2026 속리산 겨울 소원 축제: 빛의 숲, 세조의 길",
            "  - 기간: 2026년 12월 19일(토) ~ 2027년 1월 3일(토), 16일간",
            "  - 장소: 법주사, 세조길, 오리숲",
            "",
            "핵심 가치",
            "  - 치유(Healing) - 세조의 치유 스토리",
            "  - 소원(Wish) - 새해 소원 빌기",
            "  - 미디어아트(Tech) - 문화유산의 현대적 재해석"
        ],
        accent_color=BLUE
    )

    # ========== 슬라이드 7: 겨울 축제 콘텐츠 ==========
    add_content_slide(
        prs,
        "겨울 축제 킬러 콘텐츠",
        [
            "1. 법주사 팔상전 미디어아트 '미륵의 빛'",
            "   - 국보 제55호 목탑에 프로젝션 매핑",
            "   - 1막: 속리산의 사계 / 2막: 미륵의 하생 / 3막: 천년의 소원",
            "   - AR로 관람객 소원 메시지 투영",
            "",
            "2. 세조길 소원등 트레킹 '왕의 산책'",
            "   - 수천 개의 전통 등으로 빛의 터널 조성",
            "   - 나만의 소원등 만들기 + 제등 행렬",
            "   - 침묵의 묵언 트레킹 (명상 체험)",
            "",
            "3. 보은 윈터 푸드존",
            "   - 대추 라떼, 대추 찐빵, 도깨비 팥죽"
        ],
        accent_color=BLUE
    )

    # ========== 슬라이드 8: 차별화 포인트 ==========
    add_table_slide(
        prs,
        "타 축제와의 차별화",
        ["비교 항목", "속리산 축제", "합천 고스트파크", "화천 산천어축제"],
        [
            ["공간", "실제 국립공원 숲", "영화 세트장", "하천 빙판"],
            ["기후 의존도", "낮음 (빛/스토리)", "낮음", "높음 (결빙 필수)"],
            ["스토리텔링", "도깨비/세조 설화", "영화 기반", "없음"],
            ["문화유산 연계", "유네스코 법주사", "없음", "없음"],
            ["MZ 타겟팅", "AR, SNS 챌린지", "공포 체험", "체험 위주"],
        ]
    )

    # ========== 슬라이드 9: 실현가능성 (예산) ==========
    add_two_column_slide(
        prs,
        "실현가능성: 예산 계획 (총 18억 원)",
        "여름 축제 (8억 원)",
        [
            "인건비: 1.5억",
            "  - 연기자, 안전요원, 진행요원",
            "",
            "시설비: 2억",
            "  - 조명, 음향, 워터밤 시설",
            "",
            "콘텐츠 제작: 1.2억",
            "  - 시나리오, 소품, AR 앱",
            "",
            "홍보/마케팅: 1억",
            "",
            "운영비: 0.8억",
            "",
            "안전관리/예비비: 1.5억"
        ],
        "겨울 축제 (10억 원)",
        [
            "미디어아트 제작: 3억",
            "  - 프로젝션 매핑 콘텐츠",
            "",
            "조명/시설: 2.5억",
            "  - 소원등, LED, 돔텐트",
            "",
            "인건비: 1억",
            "",
            "콘텐츠 제작: 0.8억",
            "",
            "홍보/마케팅: 1억",
            "",
            "운영/안전/예비비: 1.7억"
        ]
    )

    # ========== 슬라이드 10: 추진 일정표 ==========
    add_table_slide(
        prs,
        "추진 일정표",
        ["시기", "주요 일정"],
        [
            ["2025년 12월", "공모전 당선, 기획 착수"],
            ["2026년 1~2월", "국립공원공단, 문화재청 협의"],
            ["2026년 3~4월", "콘텐츠 상세 기획, 예산 확정"],
            ["2026년 5월", "시설 설치, 리허설, 홍보 시작"],
            ["2026년 6월 27일~7월 12일", "여름 축제 개최"],
            ["2026년 7~8월", "여름 평가, 겨울 준비"],
            ["2026년 9~11월", "미디어아트 제작, 조명 설치"],
            ["2026년 12월 19일~2027년 1월 3일", "겨울 축제 개최"],
        ]
    )

    # ========== 슬라이드 11: 타겟 & 마케팅 ==========
    add_content_slide(
        prs,
        "타겟 고객 & SNS 마케팅 전략",
        [
            "타겟 고객 세분화",
            "  - MZ세대 호러매니아 (35%): 인스타그램, 틱톡, 유튜브",
            "  - 가족 단위 (30%): 네이버 블로그, 맘카페",
            "  - 연인/친구 그룹 (20%): 여행 플랫폼",
            "  - 중장년 힐링족 (15%): TV, 라디오",
            "",
            "SNS 바이럴 전략",
            "  - 해시태그: #속리산도깨비 #세조의길 #법주사빛축제",
            "  - 인플루언서 협업: 공포 유튜버, 여행 인스타그래머",
            "  - 틱톡 챌린지: #도깨비씨름챌린지",
            "",
            "인스타그래머블 포토존 4개소 설치"
        ]
    )

    # ========== 슬라이드 12: 기대효과 & 마무리 ==========
    add_table_slide(
        prs,
        "기대효과 & 심사기준 대응",
        ["심사기준", "대응 내용"],
        [
            ["독창성", "국립공원 최초 호러 트레킹, 도깨비 IP, 유네스코 미디어아트"],
            ["실현가능성", "세부 예산 18억, 협력기관 명시, 구체적 일정표"],
            ["관광객 매력도", "타겟별 맞춤 콘텐츠, SNS 전략, 포토존 설계"],
        ]
    )

    # 마지막 슬라이드에 경제효과 추가
    slide = prs.slides[-1]
    effect_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2.5))
    tf = effect_box.text_frame
    tf.word_wrap = True

    lines = [
        "예상 경제효과",
        "",
        "  총 방문객: 14.4만명 (여름 8만 + 겨울 6.4만)",
        "  직접 경제효과: 78.4억 원",
        "  간접 경제효과: 117.6억 원",
        "  총 경제효과: 196억 원",
        "",
        "\"2026년을 사계절 365일 설레는 보은 관광의 원년으로!\""
    ]

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(20)
        else:
            p = tf.add_paragraph()
            p.font.size = Pt(16)
        p.text = line
        p.font.color.rgb = DARK

    # 파일 저장
    output_path = r"C:\Users\user\songnisan-festival-2026\속리산축제_제안서.pptx"
    prs.save(output_path)
    print(f"PPT 생성 완료: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
